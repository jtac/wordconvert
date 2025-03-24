"""
PPTXCreator: Create PowerPoint presentations from generated content.
"""

import logging
import shutil
from typing import Dict, Any, Optional

import pptx

from .utils import format_slide_notes

logger = logging.getLogger(__name__)


class PPTXCreator:
    """
    Create PowerPoint presentations from AI-generated content.
    """

    def __init__(self, template_path: Optional[str], output_path: str):
        """
        Initialize the PPTX creator.

        Args:
            template_path: Path to the PowerPoint template (optional)
            output_path: Path to save the output presentation
        """
        self.template_path = template_path
        self.output_path = output_path
        self.presentation = None

    def create_presentation(
        self, presentation_structure: Dict[str, Any], slide_layouts: Dict[str, Any]
    ) -> None:
        """
        Create a PowerPoint presentation from the generated content.

        Args:
            presentation_structure: Generated presentation structure
            slide_layouts: Available slide layouts from the template
        """
        try:
            # Copy template to output file if provided, otherwise create new
            if self.template_path:
                # Create a physical copy of the template file
                shutil.copy2(self.template_path, self.output_path)

                # Now open the copy we just made
                self.presentation = pptx.Presentation(self.output_path)
                logger.info(
                    "Copied template '%s' to '%s'", self.template_path, self.output_path
                )
            else:
                # Just create a new presentation with default layouts
                self.presentation = pptx.Presentation()

            # Create slides based on content
            slides_data = presentation_structure.get("slides", [])

            for slide_data in slides_data:
                slide_type = slide_data.get("slide_type", "content")

                # Find appropriate layout
                layout = None

                # Select an appropriate layout based on slide type
                if slide_type == "title":
                    # Title slide is usually index 0
                    layout = self.presentation.slide_layouts[0]
                elif slide_type == "section":
                    # Try to find a section layout, fallback to title layout
                    if "section" in slide_layouts and slide_layouts["section"]:
                        index = slide_layouts["section"]["index"]
                        layout = self.presentation.slide_layouts[index]
                    else:
                        # Use title slide layout for sections if no specific section layout
                        layout = self.presentation.slide_layouts[0]
                else:
                    # Content slide - look for a content layout or use index 1
                    if "content" in slide_layouts and slide_layouts["content"]:
                        index = slide_layouts["content"]["index"]
                        layout = self.presentation.slide_layouts[index]
                    else:
                        # Fallback to a standard content layout
                        layouts = self.presentation.slide_layouts
                        layout = layouts[1] if len(layouts) > 1 else layouts[0]

                # Create slide with selected layout
                slide = self.presentation.slides.add_slide(layout)

                # Populate slide content based on layout type
                if slide_type == "title":
                    self._create_title_slide(slide, slide_data)
                elif slide_type == "section":
                    self._create_section_slide(slide, slide_data)
                else:  # Content or other types
                    self._create_content_slide(slide, slide_data)

            # Save the presentation
            self.presentation.save(self.output_path)
            logger.info("Saved presentation to %s", self.output_path)

        except Exception as e:
            logger.error("Error creating presentation: %s", str(e))
            raise RuntimeError(f"Failed to create presentation: {str(e)}") from e

    def _create_title_slide(
        self, slide: pptx.slide.Slide, slide_data: Dict[str, Any]
    ) -> None:
        """Create and populate a title slide."""
        # Find title and subtitle placeholders
        title_placeholder = None
        subtitle_placeholder = None

        for shape in slide.shapes:
            if not hasattr(shape, "is_placeholder") or not shape.is_placeholder:
                continue

            if hasattr(shape, "placeholder_format"):
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # Title placeholder
                    title_placeholder = shape
                elif ph_type == 2:  # Subtitle placeholder
                    subtitle_placeholder = shape

        # Set title
        if title_placeholder and "title" in slide_data:
            title_placeholder.text = slide_data["title"]

        # Set subtitle
        if subtitle_placeholder and "subtitle" in slide_data and slide_data["subtitle"]:
            subtitle_placeholder.text = slide_data["subtitle"]

        # Add notes if present
        if "notes" in slide_data and slide_data["notes"]:
            slide.notes_slide.notes_text_frame.text = format_slide_notes(
                slide_data["notes"]
            )

    def _create_section_slide(
        self, slide: pptx.slide.Slide, slide_data: Dict[str, Any]
    ) -> None:
        """Create and populate a section slide."""
        # Find title placeholder
        title_placeholder = None

        for shape in slide.shapes:
            if not hasattr(shape, "is_placeholder") or not shape.is_placeholder:
                continue

            if hasattr(shape, "placeholder_format"):
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # Title placeholder
                    title_placeholder = shape

        # Set title
        if title_placeholder and "title" in slide_data:
            title_placeholder.text = slide_data["title"]

        # Add notes if present
        if "notes" in slide_data and slide_data["notes"]:
            slide.notes_slide.notes_text_frame.text = format_slide_notes(
                slide_data["notes"]
            )

    def _create_content_slide(
        self, slide: pptx.slide.Slide, slide_data: Dict[str, Any]
    ) -> None:
        """Create and populate a content slide."""
        # Find title and content placeholders
        title_placeholder = None
        content_placeholder = None

        for shape in slide.shapes:
            if not hasattr(shape, "is_placeholder") or not shape.is_placeholder:
                continue

            if hasattr(shape, "placeholder_format"):
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # Title placeholder
                    title_placeholder = shape
                elif ph_type in (2, 3, 7, 8, 18):  # Various content placeholders
                    content_placeholder = shape

        # Set title
        if title_placeholder and "title" in slide_data:
            title_placeholder.text = slide_data["title"]

        # Set content (bullets)
        if content_placeholder and "bullets" in slide_data and slide_data["bullets"]:
            # Clear any existing text
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            # Add bullet points
            for i, bullet in enumerate(slide_data["bullets"]):
                # For first bullet, use the first paragraph
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0  # Top level bullet

        # Add notes if present
        if "notes" in slide_data and slide_data["notes"]:
            slide.notes_slide.notes_text_frame.text = format_slide_notes(
                slide_data["notes"]
            )
