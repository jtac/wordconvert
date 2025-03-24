"""
TemplateManager: Analyze and manage PowerPoint templates.
"""

import logging
from typing import Dict, List, Any

import pptx

logger = logging.getLogger(__name__)


class TemplateManager:
    """
    Analyze and manage PowerPoint templates.
    """

    # Define slide layout type mappings
    SLIDE_LAYOUT_TYPES = {
        "title": ["Title Slide", "Title"],
        "section": ["Section Header", "Section"],
        "content": ["Title and Content", "Content", "Two Content", "Comparison"],
        "picture": ["Picture with Caption", "Title and Picture"],
        "blank": ["Blank"],
    }

    def __init__(self, template_path: str):
        """
        Initialize the template manager.

        Args:
            template_path: Path to the PowerPoint template
        """
        self.template_path = template_path
        self.presentation = None

    def analyze_template(self) -> Dict[str, Any]:
        """
        Analyze the template and identify available slide layouts.

        Returns:
            Dict mapping layout types to available slide layouts
        """
        try:
            self.presentation = pptx.Presentation(self.template_path)
        except Exception as e:
            logger.error("Error opening template: %s", str(e))
            raise RuntimeError(f"Failed to open template: {str(e)}") from e

        # Map layout types to available layouts in the template
        slide_layouts = {
            "title": None,
            "section": None,
            "content": None,
            "picture": None,
            "blank": None,
        }

        # Find matching layouts
        for layout_idx, layout in enumerate(self.presentation.slide_layouts):
            layout_name = layout.name

            for layout_type, possible_names in self.SLIDE_LAYOUT_TYPES.items():
                # Check if this layout matches any of our defined types
                if any(name.lower() in layout_name.lower() for name in possible_names):
                    # If we haven't assigned this type yet, or this is a better match,
                    # update the mapping
                    if slide_layouts[layout_type] is None:
                        slide_layouts[layout_type] = {
                            "index": layout_idx,
                            "name": layout_name,
                            "placeholders": self._analyze_placeholders(layout),
                        }

        # If we couldn't find specialized layouts, use generic ones
        # Title layout is essential - if not found, use first layout
        if slide_layouts["title"] is None and len(self.presentation.slide_layouts) > 0:
            layout = self.presentation.slide_layouts[0]
            slide_layouts["title"] = {
                "index": 0,
                "name": layout.name,
                "placeholders": self._analyze_placeholders(layout),
            }

        # Content layout is essential - if not found, use second layout or first non-title
        if (
            slide_layouts["content"] is None
            and len(self.presentation.slide_layouts) > 1
        ):
            layout_idx = 1 if len(self.presentation.slide_layouts) > 1 else 0
            layout = self.presentation.slide_layouts[layout_idx]
            slide_layouts["content"] = {
                "index": layout_idx,
                "name": layout.name,
                "placeholders": self._analyze_placeholders(layout),
            }

        # Section layout - if not found, use title layout
        if slide_layouts["section"] is None and slide_layouts["title"] is not None:
            slide_layouts["section"] = slide_layouts["title"]

        return slide_layouts

    def _analyze_placeholders(self, layout) -> List[Dict[str, Any]]:
        """
        Analyze placeholders in a slide layout.

        Args:
            layout: Slide layout to analyze

        Returns:
            List of placeholder information
        """
        placeholders = []

        for placeholder in layout.placeholders:
            placeholders.append(
                {
                    "index": placeholder.placeholder_format.idx,
                    "type": placeholder.placeholder_format.type,
                    "name": placeholder.name,
                }
            )

        return placeholders
